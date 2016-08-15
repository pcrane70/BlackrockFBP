from pptx import Presentation
import math

# make presentatio

def pres():

	pres = Presentation()

	# layout
	layout_1 = pres.slide_layouts[0]

	# Slide
	slide = pres.slides.add_slide(layout_1)

	# Mek shape
	shape = slide.shapes

	# title
	title = shape.placeholders[0]

	title.text = 'Practice'

	pres.save('practice.pptx')



def testCount():

	data = 375

	for page in range(0, math.ceil( data / 4 )):

		for order in range(0, 4):

			index = (page * 4)+(order)
			print( 'Page -> ' + str( index ) + ' out of ' + str( data ) + ' -- ' + str( (index / data) * 100 ) + '%')

def testList():

	data = [1, 2, 3, 4, 5]

	for i in range(0, 5):

		print(str( len(data) ) + ' - ' + str(i))

testList()