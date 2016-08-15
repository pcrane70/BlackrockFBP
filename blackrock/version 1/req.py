#!/usr/bin/python
#coding: utf-8 
import os, sys

# Preliminary Functions

import xlrd, re, math, datetime
from pptx import Presentation
from unidecode import unidecode
from PIL import Image
from pptx.util import Inches, Pt

# Gets name of excel file from directory
def getExcelNames(directory):
	files = os.listdir(directory)
	pattern = r'(xls|xlsx|xlw){1}'

	print(' [ Status ] Identifying files in directory [ ' + str(directory) + ' ].')

	file_name = []

	for file in files:
		if re.search(pattern, file):
			add_file = {
				'name': str(file),
			}
			file_name.append(add_file)

	print(' [ Status ] Files successfully identified.')
	print( '  [ File 1 ]  {0} \n  [ File 2 ] {1}'.format( file_name[0]['name'], file_name[1]['name'] ))

	# Prompt the files that have ben selected
	print('\n \n \n ')

	identified = False
	question_index = int(input(str(file_name[0]['name']) + ' and ' + str(file_name[1]['name']) + ' have been selected. \n Please identify which excel sheet contains the questions by typing \n [ > ] 1, for the first file, and \n [ > ] 2 for the second file \n [ > ]  '))

	while not identified:

		if int(question_index) == 1:
			file_name[0]['type'] = 'question'
			file_name[1]['type'] = 'non-question'
			identified = True
			break

		elif int(question_index) == 2:
			file_name[0]['type'] = 'non-question'
			file_name[1]['type'] = 'question'
			identified = True
			break

			
		print('\n \n \n')
		print('Invalid answer. Please enter 1 for ' + str(file_name[0]['name']) + ' or 2 for ' + str(file_name[1]['name']) + '. \n')
		question_index = input('Answer: ')

	return file_name

# Get contents from excel file
def getExcelContents(excel_file_name, question_file_name):
	
	# Step 1:  Get data from excel file name

	workbook = xlrd.open_workbook(excel_file_name, on_demand=True, encoding_override="cp1252")
	worksheet = workbook.sheet_by_index(0)

	workbook2 = xlrd.open_workbook(question_file_name, on_demand=True, encoding_override="cp1252")
	worksheet2 = workbook2.sheet_by_index(0)

	# Empty List
	info = []

	for i in range(1, 1500):
		if worksheet.cell(i, 0).value == '' or worksheet.cell(i, 0).value == None:
			break

		else:
			data = {
				'first_name': worksheet.cell_value(i, 0),
				'last_name': worksheet.cell_value(i, 1),
				'organization': worksheet.cell_value(i, 3),
				'location': worksheet.cell_value(i, 29),
				'university': worksheet.cell_value(i, 19),
			}

			if data['organization'] == 42:
				data['organization'] = worksheet.cell_value(i, 9)

			# Check if the name is equal to the name on the question
			for x in range(1, 376):
				if data['first_name'] == worksheet2.cell(x, 2).value and data['last_name'] == worksheet2.cell(x, 1).value:
					

					interests = worksheet2.cell_value(x, 22)

					dinner = worksheet2.cell_value(x, 23)
					
					film = worksheet2.cell_value(x, 24)

					# Add answers to question
					data['questions'] = {
						'interest': str(interests),
						'dinner': str(dinner),
						'film': str(film),
					}

			# Add the data to info list
			info.append(data)

	# Step 2: Get data from question file name

	return info

# Get image contents
def getImageContents(image_directory, excel_cont):

	# Excel contents should return a list of dictionaries
	for person in excel_cont:
		# Each dictionary

		name = person['last_name'] + ', ' + person['last_name']
		
		photo_dir = os.listdir(image_directory)
		profiles = []

		pattern = r'(\.){1}(jpeg|jpg|gif|jfif|tiff|png){1}$'

		# Check if it is an image
		for file in photo_dir:
			if re.search(pattern, file):

				# Add to profiles list
				profiles.append(file)

	return profiles


def combineImageExcel(image_list, excel_list):

	# Make pattern to get name
	lastpatt = r'(?P<last>[a-zA-Z]+)(\,){1}(\s)?(?P<first>[a-zA-Z]+)(\s)?(\.){1}(?P<ext>[a-zA-Z]+)'
	firstpatt = r'(?P<first>[a-zA-Z]+)(\s)+(?P<last>[a-zA-Z]+)(\.){1}(?P<ext>[a-zA-Z]+)'

	for img in image_list:

		if re.match(lastpatt, img):
			last = re.match(lastpatt, img).group('last')
			first = re.match(lastpatt, img).group('first')
			ext = re.match(lastpatt, img).group('ext')
		
		elif re.match(firstpatt, img):
			last = re.match(firstpatt, img).group('last')
			first = re.match(firstpatt, img).group('first')
			ext = re.match(firstpatt, img).group('ext')

		# Check if image name matches that in excel_list
		for person in excel_list:

			first_name = person['first_name']
			last_name = person['last_name']

			if first_name == first and last_name == last:

				person['image'] = str(img)

	return excel_list

# Turn images given to reasonably sized images
def imgToThumb(directory, data):

	# Get image name form dict data

	thumb_size = 100, 100

	for person in data:

		# Get file location
		if 'image' in person:
			file = directory + '/' + person['image']

		# Make Dir

			if not os.path.exists(directory + '/thumb'):
				os.mkdir(directory + '/thumb')
			
			thumb_file = directory + '/thumb/thumb - ' + person['image']

			# Get the image
			image = Image.open(file)

			# Create thumbnail
			image.thumbnail(thumb_size)

			# Save Image
			image.convert('RGB').save( thumb_file , "JPEG")

			# Status Update
			print(person['image'] + ' updated.')

			# Update data
			person['thumb'] = thumb_file

	return data



# Create presentation
def makePres(data, directory):

	all_groups = []

	# Split data by organization
	for person in data:

		all_groups.append( person['organization'] )

	# Make all_groups into a set
	discrete = set()

	for org in all_groups:
		discrete.add(org)

	# Step 1 - Make a main presentation
	# Step 2 - Make a presentation for each group

	year = datetime.date.today().year

	# [1] Main presentation
	main_pres = 'Blackrock Facebook ' + str( year )
	main_data = data

	pres(main_data, main_pres, directory)

	# [2] Separate Presentation
	for org in discrete:
		org_name = str(year) + ' ' + str(org) + ' Facebook'
		org_data = []

		for person in data:
			if person['organization'] == org:
				org_data.append(person)

		pres(org_data, org_name, directory)

	print(' Your Blackrock Facebook Profiler is complete. ')

# General presentation
def pres(data, name, directory):
	
	pres = Presentation()

	# Set slide dimentions
	pres.slide_width = 6858000
	pres.slide_height = 9144000

	# Get title slide
	title_layout = pres.slide_layouts[0]

	# Get first slide
	title_slide = pres.slides.add_slide( title_layout )

	title = title_slide.shapes.title
	title.text = name

	# Make each page
	for page in range( 0, math.ceil( len(data)/4 ) ):

		# Make slide
		page_layout = pres.slide_layouts[1]
		slide = pres.slides.add_slide( page_layout )

		# Get shape
		shape = slide.shapes

		for order in range(0, 4):

			position = (page * 4) + (order)
			# Make sure it is in the data
			if not position > (len(data) - 1):
				# Make sure in range of data list
				
				# Get positions
				left = Inches(3)
				width = Inches(5)
				height = Inches(2)

				# Get top Position
				if order == 0:
					top_inches = 1
					top_par = 1.2
				else:
					top_inches = (order * 2) + 1
					top_par = 1.2

				top = Inches( top_inches )
				top_paragraph = Inches( top_par )

				# Create Image
				if 'image' in data[position]:
					image_path = str(directory) + '/Photos/' + str(data[position]['image'])
					image_left = Inches(1)

					pic = shape.add_picture(image_path, image_left, top)

			# Textbox
				txBox = shape.add_textbox(left, top, width, height)
				txBox2 = shape.add_textbox(left, top_paragraph, width, height)
				
				name = txBox.text_frame
				name.text = str( data[position]['first_name'] ) + ' ' + str( data[position]['last_name'] )

				content = txBox2.text_frame

				if 'questions' in data[position]:	
					content.text = 'University: {0} \nInterest: {1} \nFavorite Film: {2} \nDinner Invites: {3}'.format( str(data[position]['university']), str(data[position]['questions']),  str(data[position]['questions']),  str(data[position]['questions']) )
				else:
					content.text = ''


	# Save presentation
	pres_name = str(name) + '.pptx'
	pres.save( pres_name )