from pptx import Presentation as Pres

# New Presentation
prs = Pres()

# Set Dimentions
prs.slide_width = 6858000
prs.slide_height = 9144000

# Title Slide
title_slide_layout = prs.slide_layouts[0]
# Add Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]



title.text = 'Test'
subtitle.text = '123'

prs.save('work.pptx')

# 14 Powerpoints - 13 Split, 1 Final
# Get 